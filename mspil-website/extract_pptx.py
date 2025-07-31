#!/usr/bin/env python3
"""Extract text content from PowerPoint presentation"""

import sys
from pptx import Presentation

def extract_text_from_pptx(file_path):
    """Extract all text from a PowerPoint file"""
    try:
        prs = Presentation(file_path)
        
        all_text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            slide_text.append(f"\n{'='*60}")
            slide_text.append(f"SLIDE {slide_num}")
            slide_text.append(f"{'='*60}")
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
                
                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_text.append(cell_text)
                        if row_text:
                            slide_text.append(" | ".join(row_text))
            
            if len(slide_text) > 3:  # More than just the header
                all_text.extend(slide_text)
        
        return "\n".join(all_text)
        
    except Exception as e:
        return f"Error reading PowerPoint file: {str(e)}"

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python extract_pptx.py <path_to_pptx_file>")
        sys.exit(1)
    
    file_path = sys.argv[1]
    extracted_text = extract_text_from_pptx(file_path)
    print(extracted_text)